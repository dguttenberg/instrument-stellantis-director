"""Deck extraction → draft Script (fake Claude), and the retrieval-vs-generative
prompt guidance that keeps TwelveLabs queries broad."""

import pytest

from director_agent.ingest.extractor import DeckExtractor
from director_agent.director.prompt import build_system_prompt
from director_agent.schemas.dials import Dials, StylingInputs

from conftest import FakeAnthropic

EXTRACT = {
    "scenes": [
        {"scene_index": 1, "scene_description": "Ram drives toward camera",
         "super_called": True, "super_intent": "Giving around here.",
         "cell_type": "regionalized_running_w_cgi_ai", "suggested_trim": "D28H91"},
        {"scene_index": 2, "scene_description": "Stock establishing town",
         "super_called": False, "super_intent": None, "cell_type": "stock", "suggested_trim": None},
        {"scene_index": 3, "scene_description": "HEMI badge close-up",
         "super_called": True, "super_intent": "Strongest Ram ever.",
         "cell_type": "existing_running_footage", "suggested_trim": None},
    ]
}


def _extractor():
    return DeckExtractor(client=FakeAnthropic(EXTRACT, tool_name="emit_draft_script"))


def test_extract_from_text_builds_script():
    script = _extractor().extract_from_text("RAM Holiday Deck.pptx", "deck text")
    assert script.script_id == "ram_holiday_deck"
    assert script.total_scenes == 3
    assert [s.cell_type for s in script.scenes] == [
        "regionalized_running_w_cgi_ai", "stock", "existing_running_footage"]
    assert script.scenes[0].cell_id == "ram_holiday_deck_scene1"
    assert script.scenes[0].trim_intent == "D28H91"
    # regionalized_ai gets default camera angles; others empty
    assert script.scenes[2].camera_angles == []


def test_extract_rejects_unknown_type():
    with pytest.raises(ValueError):
        _extractor().extract("notes.txt", b"hi")


def test_pptx_text_extraction():
    """Exercise the real python-pptx path (shapes + table + notes) on a synthetic deck."""
    import io

    from pptx import Presentation
    from pptx.util import Inches

    from director_agent.ingest.pptx_text import slides_text

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Scene 1"
    slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1)).text_frame.text = (
        "Ram drives toward camera"
    )
    table = slide.shapes.add_table(1, 2, Inches(1), Inches(3), Inches(5), Inches(1)).table
    table.cell(0, 0).text = "technique"
    table.cell(0, 1).text = "running footage"
    slide.notes_slide.notes_text_frame.text = "super: The holidays are all about giving."

    buf = io.BytesIO()
    prs.save(buf)
    texts = slides_text(buf.getvalue())
    assert len(texts) == 1
    body = texts[0]
    assert "Scene 1" in body
    assert "Ram drives toward camera" in body
    assert "technique | running footage" in body
    assert "[notes] super: The holidays" in body


def test_retrieval_vs_generative_guidance_present():
    sp = build_system_prompt(StylingInputs(tone_of_voice="x", creative_intent="x", direction="x",
                                           lighting_aesthetic="x", storytelling_flow="x"),
                             Dials(), "existing_running_footage")
    assert "RETRIEVAL queries and GENERATIVE prompts" in sp
    assert "15 words" in sp                       # findable cap (re-tuned up from 12)
    assert "broad but regionally flavored" in sp  # light location intelligence retained
    assert "PRIMARY/preferred" in sp              # core-footage cell guidance headline
