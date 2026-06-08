"""Extract readable text per slide from a .pptx (shapes, tables, speaker notes)."""

from __future__ import annotations

import io
from typing import List


def slides_text(file_bytes: bytes) -> List[str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    out: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[notes] {notes}")
        out.append(f"--- Slide {i} ---\n" + "\n".join(parts))
    return out
